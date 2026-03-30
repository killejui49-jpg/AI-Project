#!/usr/bin/env python3
"""Generate AI Presentation PPTX for Tronc Commun Scientifique."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0f, 0x0c, 0x29)
BG_MID = RGBColor(0x1a, 0x1a, 0x2e)
PURPLE = RGBColor(0x81, 0x8c, 0xf8)
LIGHT_PURPLE = RGBColor(0xc0, 0x84, 0xfc)
PINK = RGBColor(0xf4, 0x72, 0xb6)
WHITE = RGBColor(0xff, 0xff, 0xff)
WHITE_70 = RGBColor(0xb3, 0xb3, 0xb3)
WHITE_50 = RGBColor(0x80, 0x80, 0x80)
GREEN = RGBColor(0x34, 0xd3, 0x99)
RED = RGBColor(0xf8, 0x71, 0x71)
CARD_BG = RGBColor(0x1e, 0x1b, 0x4b)


def add_bg(slide, color=BG_DARK):
    """Add solid background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=0):
    """Add a rounded rectangle shape as background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box with given properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(6)):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


# =================== SLIDE 1: TITLE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# Decorative accent bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

# Robot emoji / icon text
add_text_box(slide, Inches(5.5), Inches(1.0), Inches(2.5), Inches(1.2),
             "AI", font_size=72, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.5),
             "L'Intelligence Artificielle", font_size=48, color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.8),
             "Projet Scolaire - Informatique", font_size=24,
             color=WHITE_70, alignment=PP_ALIGN.CENTER)

# Student info box
box = add_shape_bg(slide, Inches(4), Inches(4.8), Inches(5.3), Inches(1.8), CARD_BG, 0.05)
add_text_box(slide, Inches(4.2), Inches(4.9), Inches(5), Inches(0.5),
             "Realise par", font_size=16, color=WHITE_70, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(5.3), Inches(5), Inches(0.6),
             "Ayoub Assouar", font_size=28, color=LIGHT_PURPLE,
             bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(5.9), Inches(5), Inches(0.5),
             "Tronc Commun Scientifique", font_size=16,
             color=WHITE_50, alignment=PP_ALIGN.CENTER)

# Bottom accent
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), LIGHT_PURPLE)


# =================== SLIDE 2: SOMMAIRE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
             "Sommaire", font_size=40, color=PURPLE, bold=True)

chapters = [
    "1.  Qu'est-ce que l'IA ?",
    "2.  Histoire de l'IA",
    "3.  Les types d'IA",
    "4.  Comment fonctionne l'IA ?",
    "5.  Applications de l'IA",
    "6.  L'IA dans la vie quotidienne",
    "7.  Avantages et Inconvenients",
    "8.  L'avenir de l'IA",
]

for i, ch in enumerate(chapters):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(1.25)
    box = add_shape_bg(slide, x, y, Inches(5.8), Inches(1.0), CARD_BG, 0.03)
    # Number circle
    num_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6)
    )
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = PURPLE
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    # Text
    add_text_box(slide, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.6),
                 ch.split(". ", 1)[1] if ". " in ch else ch,
                 font_size=18, color=WHITE_70)


# =================== SLIDE 3: Qu'est-ce que l'IA ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 1", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Qu'est-ce que l'Intelligence Artificielle ?",
             font_size=36, color=WHITE, bold=True)

# Left text
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(7), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "L'intelligence artificielle (IA) est une branche de l'informatique qui vise a creer des systemes capables d'effectuer des taches qui necessitent normalement l'intelligence humaine."
p.font.size = Pt(17)
p.font.color.rgb = WHITE_70
p.font.name = 'Calibri'
p.space_after = Pt(14)

add_paragraph(tf, "Ces taches incluent : la reconnaissance vocale, la prise de decision, la traduction de langues, et la perception visuelle.",
              font_size=17, color=WHITE_70)

add_paragraph(tf, 'Le terme "Intelligence Artificielle" a ete invente par John McCarthy en 1956 lors de la conference de Dartmouth.',
              font_size=17, color=WHITE_70)

add_paragraph(tf, "L'IA cherche a simuler les processus cognitifs humains comme l'apprentissage, le raisonnement, et l'auto-correction.",
              font_size=17, color=WHITE_70)

# Right visual - brain icon box
box = add_shape_bg(slide, Inches(8.5), Inches(2.2), Inches(4), Inches(4), CARD_BG, 0.05)
add_text_box(slide, Inches(8.5), Inches(3.0), Inches(4), Inches(1.5),
             "IA", font_size=80, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.5), Inches(4.5), Inches(4), Inches(0.5),
             "Intelligence Artificielle", font_size=14, color=WHITE_50, alignment=PP_ALIGN.CENTER)


# =================== SLIDE 4: Histoire ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 2", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'Histoire de l'Intelligence Artificielle",
             font_size=36, color=WHITE, bold=True)

timeline = [
    ("1950", "Alan Turing propose le 'Test de Turing' pour mesurer l'intelligence des machines"),
    ("1956", "Conference de Dartmouth : naissance officielle de l'IA comme domaine de recherche"),
    ("1966", "ELIZA, le premier chatbot, est cree par Joseph Weizenbaum au MIT"),
    ("1997", "Deep Blue (IBM) bat le champion du monde d'echecs Garry Kasparov"),
    ("2011", "IBM Watson gagne au jeu televise Jeopardy! contre des champions humains"),
    ("2016", "AlphaGo (Google DeepMind) bat le champion mondial du jeu de Go"),
    ("2022+", "ChatGPT et les grands modeles de langage revolutionnent l'IA generative"),
]

# Timeline line
add_shape_bg(slide, Inches(1.5), Inches(1.9), Inches(0.06), Inches(5.2), PURPLE)

for i, (year, desc) in enumerate(timeline):
    y = Inches(1.9) + i * Inches(0.74)
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), y, Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = LIGHT_PURPLE
    dot.line.fill.background()
    # Year
    add_text_box(slide, Inches(2.1), y - Inches(0.02), Inches(1.2), Inches(0.4),
                 year, font_size=16, color=LIGHT_PURPLE, bold=True)
    # Description
    add_text_box(slide, Inches(3.3), y - Inches(0.02), Inches(9), Inches(0.4),
                 desc, font_size=15, color=WHITE_70)


# =================== SLIDE 5: Types d'IA ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 3", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Les Types d'Intelligence Artificielle",
             font_size=36, color=WHITE, bold=True)

types = [
    ("IA Faible (Narrow AI)", "Concue pour une tache specifique. C'est le type d'IA le plus courant aujourd'hui.\n\nExemples : Siri, Alexa, filtres de spam, recommandations Netflix."),
    ("IA Forte (General AI)", "Une IA qui pourrait comprendre et apprendre n'importe quelle tache intellectuelle comme un humain.\n\nCe type n'existe pas encore mais fait l'objet de recherches actives."),
    ("Super IA (Super AI)", "Une IA hypothetique qui depasserait l'intelligence humaine dans tous les domaines.\n\nC'est un concept theorique qui souleve beaucoup de questions ethiques."),
]

icons = ["Narrow", "General", "Super"]
colors = [PURPLE, LIGHT_PURPLE, PINK]

for i, (title, desc) in enumerate(types):
    x = Inches(0.6) + i * Inches(4.2)
    box = add_shape_bg(slide, x, Inches(2.0), Inches(3.9), Inches(4.8), CARD_BG, 0.04)
    # Icon area
    icon_box = add_shape_bg(slide, x + Inches(1.2), Inches(2.3), Inches(1.5), Inches(1.5), colors[i], 0.1)
    add_text_box(slide, x + Inches(1.2), Inches(2.5), Inches(1.5), Inches(1.2),
                 icons[i][0], font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.3), Inches(4.0), Inches(3.3), Inches(0.6),
                 title, font_size=18, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.3), Inches(4.6), Inches(3.3), Inches(2.0),
                 desc, font_size=13, color=WHITE_70, alignment=PP_ALIGN.CENTER)


# =================== SLIDE 6: Comment ca marche ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 4", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Comment Fonctionne l'IA ?",
             font_size=36, color=WHITE, bold=True)

how_items = [
    ("1. Les Donnees", "L'IA a besoin de grandes quantites de donnees pour apprendre. Plus les donnees sont nombreuses et de qualite, meilleure sera l'IA."),
    ("2. Les Algorithmes", "Des regles mathematiques permettent a l'IA d'analyser les donnees, de trouver des patterns et de prendre des decisions."),
    ("3. L'Apprentissage", "L'IA s'ameliore en s'entrainant sur des donnees. Elle ajuste ses parametres pour obtenir de meilleurs resultats a chaque iteration."),
    ("4. Reseaux de Neurones", "Inspires du cerveau humain, ces reseaux sont composes de couches de 'neurones' artificiels qui traitent l'information de maniere hierarchique."),
]

for i, (title, desc) in enumerate(how_items):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.4)
    y = Inches(2.0) + row * Inches(2.6)
    box = add_shape_bg(slide, x, y, Inches(6.0), Inches(2.3), CARD_BG, 0.04)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.3), Inches(5.2), Inches(0.5),
                 title, font_size=20, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.9), Inches(5.2), Inches(1.2),
                 desc, font_size=15, color=WHITE_70)


# =================== SLIDE 7: Applications ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 5", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Les Applications de l'IA",
             font_size=36, color=WHITE, bold=True)

apps = [
    ("Medecine", "Diagnostic de maladies, analyse d'imagerie medicale, decouverte de medicaments, chirurgie assistee par robot."),
    ("Transport", "Voitures autonomes (Tesla, Waymo), optimisation du trafic, systemes de navigation intelligents."),
    ("Education", "Tutorat personnalise, correction automatique, plateformes d'apprentissage adaptatif."),
    ("Jeux Video", "Personnages non-joueurs (PNJ) intelligents, generation procedurale de contenu, matchmaking."),
    ("Finance", "Detection de fraudes, trading algorithmique, evaluation des risques, chatbots de service client."),
    ("Environnement", "Prevision climatique, surveillance de la deforestation, optimisation energetique, agriculture de precision."),
]

for i, (title, desc) in enumerate(apps):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(2.0) + row * Inches(2.6)
    box = add_shape_bg(slide, x, y, Inches(3.9), Inches(2.3), CARD_BG, 0.04)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.5),
                 title, font_size=20, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(1.2),
                 desc, font_size=14, color=WHITE_70)


# =================== SLIDE 8: Vie quotidienne ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 6", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'IA dans Notre Vie Quotidienne",
             font_size=36, color=WHITE, bold=True)

daily = [
    ("Assistants vocaux", "Siri, Google Assistant, Alexa comprennent et repondent a vos questions"),
    ("Reseaux sociaux", "L'IA selectionne le contenu de votre fil d'actualite (Instagram, TikTok, YouTube)"),
    ("Streaming", "Netflix et Spotify utilisent l'IA pour recommander des films et de la musique"),
    ("Photographie", "Le mode portrait et les filtres de votre telephone utilisent l'IA"),
    ("Traduction", "Google Translate utilise l'IA pour traduire plus de 100 langues"),
    ("ChatGPT", "Des millions de personnes utilisent l'IA generative pour ecrire, coder et creer"),
]

for i, (title, desc) in enumerate(daily):
    y = Inches(2.0) + i * Inches(0.85)
    box = add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.72), CARD_BG, 0.02)
    # Bullet
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.22), Inches(0.28), Inches(0.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(1.7), y + Inches(0.15), Inches(2.5), Inches(0.45),
                 title, font_size=16, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(4.2), y + Inches(0.15), Inches(8), Inches(0.45),
                 desc, font_size=15, color=WHITE_70)


# =================== SLIDE 9: Avantages & Inconvenients ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 7", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Avantages et Inconvenients de l'IA",
             font_size=36, color=WHITE, bold=True)

# Advantages box
add_shape_bg(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(5.0), CARD_BG, 0.04)
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(0.6),
             "Avantages", font_size=24, color=GREEN, bold=True)

pros = [
    "Automatisation des taches repetitives et ennuyeuses",
    "Analyse rapide de grandes quantites de donnees",
    "Disponible 24h/24, 7j/7 sans fatigue",
    "Amelioration des diagnostics medicaux",
    "Reduction des erreurs humaines",
    "Innovation dans tous les domaines scientifiques",
]
for i, pro in enumerate(pros):
    y = Inches(3.0) + i * Inches(0.6)
    add_text_box(slide, Inches(1.0), y, Inches(5.2), Inches(0.5),
                 "+ " + pro, font_size=14, color=WHITE_70)

# Disadvantages box
add_shape_bg(slide, Inches(6.9), Inches(2.0), Inches(5.8), Inches(5.0), CARD_BG, 0.04)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.6),
             "Inconvenients", font_size=24, color=RED, bold=True)

cons = [
    "Risque de perte d'emplois (chomage technologique)",
    "Questions de vie privee et surveillance",
    "Biais dans les algorithmes et discrimination",
    "Dependance excessive a la technologie",
    "Cout eleve de developpement",
    "Risques de mauvaise utilisation (deepfakes, armes)",
]
for i, con in enumerate(cons):
    y = Inches(3.0) + i * Inches(0.6)
    add_text_box(slide, Inches(7.3), y, Inches(5.2), Inches(0.5),
                 "- " + con, font_size=14, color=WHITE_70)


# =================== SLIDE 10: L'avenir ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 8", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'Avenir de l'Intelligence Artificielle",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(0.6),
             "L'IA continue d'evoluer a une vitesse impressionnante. Voici quelques tendances :",
             font_size=17, color=WHITE_70)

future = [
    ("IA Generative", "Creation de textes, images, videos et musique de plus en plus realistes"),
    ("Medecine personnalisee", "Traitements adaptes a chaque patient grace a l'IA"),
    ("Robots intelligents", "Des robots capables d'interagir naturellement avec les humains"),
    ("IA ethique", "Developpement de lois et reglementations pour encadrer l'IA"),
]

for i, (title, desc) in enumerate(future):
    y = Inches(2.7) + i * Inches(1.1)
    box = add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.9), CARD_BG, 0.03)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.3), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PINK
    dot.line.fill.background()
    add_text_box(slide, Inches(1.8), y + Inches(0.15), Inches(3), Inches(0.5),
                 title, font_size=18, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(4.8), y + Inches(0.2), Inches(7.5), Inches(0.5),
                 desc, font_size=16, color=WHITE_70)


# =================== SLIDE 11: Conclusion ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CONCLUSION", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(1.0),
             "L'IA : Un Outil Puissant a Utiliser\navec Responsabilite",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(2.5), Inches(9.3), Inches(1.2),
             "L'intelligence artificielle est l'une des technologies les plus transformatrices de notre epoque. Elle offre des possibilites enormes pour ameliorer notre vie, mais elle souleve aussi des defis importants.",
             font_size=17, color=WHITE_70, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(1.2),
             "En tant qu'etudiants et futurs citoyens, il est essentiel de comprendre comment l'IA fonctionne, de connaitre ses avantages et ses limites, et de contribuer a son developpement de maniere ethique et responsable.",
             font_size=17, color=WHITE_70, alignment=PP_ALIGN.CENTER)

# Quote box
quote_box = add_shape_bg(slide, Inches(3), Inches(5.3), Inches(7.3), Inches(1.5), CARD_BG, 0.04)
add_text_box(slide, Inches(3.3), Inches(5.5), Inches(6.7), Inches(0.7),
             '"L\'intelligence artificielle est la nouvelle electricite."',
             font_size=20, color=LIGHT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.3), Inches(6.2), Inches(6.7), Inches(0.4),
             "- Andrew Ng, expert en IA",
             font_size=14, color=WHITE_50, alignment=PP_ALIGN.CENTER)


# =================== SLIDE 12: Merci ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5),
             "Merci !", font_size=64, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.8),
             "Presentation realisee par Ayoub Assouar",
             font_size=22, color=WHITE_70, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
             "Tronc Commun Scientifique",
             font_size=18, color=WHITE_50, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
             "Des questions ?",
             font_size=20, color=WHITE_50, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GREEN)


# Save
prs.save('AI_Presentation_Ayoub_Assouar.pptx')
print("Presentation saved: AI_Presentation_Ayoub_Assouar.pptx")
